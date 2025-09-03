from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from PIL import Image, ImageOps
import io, math
from typing import Dict
from app.core.processing import Grupo, Foto
from app.utils.nlg_utils import agrupa_y_redacta

def _add_textbox(slide, left, top, width, height, text, size=11):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    return tb

def export_groups_to_pptx_report(grupos: Dict[str, Grupo], archivos: Dict[str, bytes],
                                 output_pptx_path: str, max_px: int = 1600, progress_callback=None) -> None:
    prs = Presentation()
    prs.slide_width = Inches(8.27)
    prs.slide_height = Inches(11.69)
    blank_layout = prs.slide_layouts[6]

    cols, rows = 4, 3
    margin_x = 0.4
    margin_y_top = 1.2
    margin_y_bottom = 0.4
    enumerated_h = 2.5
    spacing_x = 0.1
    spacing_y = 0.1

    slide_w_in = prs.slide_width / 914400.0
    slide_h_in = prs.slide_height / 914400.0
    photo_area_h = slide_h_in - margin_y_top - enumerated_h - margin_y_bottom
    cell_w = (slide_w_in - 2 * margin_x - (cols - 1) * spacing_x) / cols
    cell_h = (photo_area_h - (rows - 1) * spacing_y) / rows

    # Calcular el número total de diapositivas para el progreso
    total_slides = sum(math.ceil(len(g.fotos) / (cols * rows)) for g in grupos.values() if g.fotos)
    slides_done = 0

    for gname, grupo in grupos.items():
        per_slide = cols * rows
        total = len(grupo.fotos)
        pages = math.ceil(total / per_slide) if per_slide else 0

        for page in range(pages):
            slide = prs.slides.add_slide(blank_layout)

            # --- INICIO DE LA CORRECCIÓN AVANZADA DE TÍTULO ---
            title_box = slide.shapes.add_textbox(
                Inches(margin_x), Inches(0.3),
                Inches(slide_w_in - 2 * margin_x), Inches(0.75) # Mantener altura aumentada
            )
            tf = title_box.text_frame
            # 1. Activar autoajuste para que el texto se reduzca para caber
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            # 2. Asegurar que el texto se divida en varias líneas
            tf.word_wrap = True
            
            p_title = tf.paragraphs[0]
            p_title.text = gname
            # 3. Establecer un tamaño de fuente INICIAL grande. Autoajuste lo reducirá si es necesario.
            p_title.font.size = Pt(16)
            # --- FIN DE LA CORRECCIÓN AVANZADA DE TÍTULO ---

            label_y = margin_y_top - 0.3
            label_box = slide.shapes.add_textbox(
                Inches(margin_x), Inches(label_y),
                Inches(3), Inches(0.4)
            )
            p_label = label_box.text_frame.paragraphs[0]
            p_label.text = "FOTOGRAFÍAS:"
            if p_label.runs:
                run_label = p_label.runs[0]
                run_label.font.bold = True
                run_label.font.size = Pt(12)

            chunk = grupo.fotos[page * per_slide:(page + 1) * per_slide]
            for idx, foto in enumerate(chunk):
                r = idx // cols
                c = idx % cols
                x = margin_x + c * (cell_w + spacing_x)
                y = margin_y_top + r * (cell_h + spacing_y)

                # Intentar diferentes variaciones de la ruta
                posible_paths = [
                    f"{foto.carpeta}/{foto.filename}",  # Unix style
                    f"{foto.carpeta}\\{foto.filename}",  # Windows style
                    foto.filename,  # Solo nombre del archivo
                    f"{foto.carpeta.replace('/', '\\')}\\{foto.filename}"  # Windows style con carpeta normalizada
                ]
                
                img_data = None
                for path in posible_paths:
                    img_data = archivos.get(path)
                    if img_data:
                        break

                if img_data:
                    try:
                        # Abrir y procesar la imagen
                        img = Image.open(io.BytesIO(img_data))
                        img = ImageOps.exif_transpose(img)
                        
                        # Asegurarse de que la imagen esté en el modo correcto
                        if img.mode in ("RGBA", "LA", "P"):
                            img = img.convert("RGB")
                            
                        # Redimensionar la imagen manteniendo la proporción
                        img.thumbnail((max_px, max_px), Image.LANCZOS)
                        
                        # Guardar en buffer
                        buffer = io.BytesIO()
                        img.save(buffer, format="JPEG", quality=80, optimize=True)
                        buffer.seek(0)
                        
                        # Agregar al slide
                        slide.shapes.add_picture(
                            buffer, Inches(x), Inches(y),
                            width=Inches(cell_w), height=Inches(cell_h)
                        )
                    except Exception as e:
                        print(f"Error procesando imagen {foto.filename}: {str(e)}")

            enum_y = slide_h_in - enumerated_h - margin_y_bottom
            enum_x = margin_x
            enum_w = slide_w_in - 2 * margin_x
            header_h = 0.4
            content_h = enumerated_h - header_h
            details_w = enum_w * 0.7
            recom_w = enum_w - details_w

            header_det = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(enum_x), Inches(enum_y),
                Inches(details_w), Inches(header_h)
            )
            header_det.fill.solid()
            header_det.fill.fore_color.rgb = RGBColor(217, 217, 217)
            header_det.line.color.rgb = header_det.fill.fore_color.rgb
            txt_hd = header_det.text_frame
            txt_hd.text = "UBICACIÓN Y DETALLE:"
            if txt_hd.paragraphs[0].runs:
                txt_hd.paragraphs[0].runs[0].font.bold = True

            header_rec = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(enum_x + details_w), Inches(enum_y),
                Inches(recom_w), Inches(header_h)
            )
            header_rec.fill.solid()
            header_rec.fill.fore_color.rgb = RGBColor(217, 217, 217)
            header_rec.line.color.rgb = header_rec.fill.fore_color.rgb
            txt_hr = header_rec.text_frame
            txt_hr.text = "RECOMENDACIONES:"
            if txt_hr.paragraphs[0].runs:
                txt_hr.paragraphs[0].runs[0].font.bold = True

            body_det = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(enum_x), Inches(enum_y + header_h),
                Inches(details_w), Inches(content_h)
            )
            body_det.fill.solid()
            body_det.fill.fore_color.rgb = RGBColor(255, 255, 255)
            body_det.line.color.rgb = body_det.fill.fore_color.rgb
            tf_det = body_det.text_frame
            tf_det.clear()

            # 1) Construir entradas (descripcion, variable) desde tus dos textos
            entradas = []
            for foto in chunk:
                full_detail = foto.specific_detail
                detail_after_plus = full_detail.split('+', 1)[1].strip() if '+' in full_detail else full_detail
                entradas.append((detail_after_plus, foto.carpeta))

            # 2) Generar oraciones agrupadas con NLG (umbral ajustable)
            oraciones = agrupa_y_redacta(entradas, umbral_similitud=0.8)

            # 3) Pintar el resultado enumerado en la caja de “UBICACIÓN Y DETALLE”
            for idx, sentencia in enumerate(oraciones, start=1):
                p = tf_det.add_paragraph()
                p.text = f"{idx}. {sentencia}"
                if p.runs:
                    run = p.runs[0]
                    run.font.size = Pt(10)
                    run.font.color.rgb = RGBColor(0, 0, 0)

            body_rec = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(enum_x + details_w), Inches(enum_y + header_h),
                Inches(recom_w), Inches(content_h)
            )
            body_rec.fill.solid()
            body_rec.fill.fore_color.rgb = RGBColor(226, 240, 217)
            body_rec.line.color.rgb = body_rec.fill.fore_color.rgb
            tf_rec = body_rec.text_frame
            tf_rec.clear()

            recs = getattr(grupo, "recomendaciones", None) or []
            rec_text = "\n".join(f"• {r}" for r in recs) if recs else "—"

            p = tf_rec.paragraphs[0]
            p.text = rec_text
            if p.runs:
                p.runs[0].font.size = Pt(10)
            
            slides_done += 1
            if progress_callback:
                progress_percentage = int((slides_done / total_slides) * 100) if total_slides > 0 else 0
                progress_callback.emit(progress_percentage)

    prs.save(output_pptx_path)

